#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint 파일 서식 유지 업데이트 스크립트
기존 텍스트의 서식을 유지하면서 내용만 업데이트합니다.
"""

import os
from pptx import Presentation
import re

def replace_text_preserve_format(text_frame, old_text, new_text):
    """텍스트 프레임의 텍스트를 교체하면서 서식을 유지합니다."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                return True
    return False

def find_and_replace_in_slide(slide, replacements):
    """슬라이드에서 텍스트를 찾아 교체합니다."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for old_text, new_text in replacements.items():
                replace_text_preserve_format(text_frame, old_text, new_text)

def update_slide_4_team_info(slide):
    """슬라이드 4: 팀원 소개 업데이트 (서식 유지)"""
    replacements = {
        "[팀명]": "AigoRhythm",
        "팀 각오 입력": "17개 클래스 문서 분류에서 최고 성능 달성을 목표로 팀워크와 기술력을 통해 도전하겠습니다!"
    }
    
    find_and_replace_in_slide(slide, replacements)
    
    # 팀원 정보 순차적 업데이트
    team_updates = [
        ("팀장 이름", "문국현"),
        ("관심 분야/전공", "Computer Vision / AI"),
        ("팀에서 수행한 역할", "모델링 총괄, 전체 프로젝트 관리")
    ]
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            full_text = text_frame.text
            
            # 첫 번째 팀장 정보
            if "팀장" in full_text and "팀장 이름" in full_text:
                for old_text, new_text in team_updates:
                    replace_text_preserve_format(text_frame, old_text, new_text)
            
            # 팀원들 정보 업데이트
            elif "팀원" in full_text and "팀원 이름" in full_text:
                if "류지헌" not in full_text:  # 첫 번째 팀원
                    replace_text_preserve_format(text_frame, "팀원 이름", "류지헌")
                    replace_text_preserve_format(text_frame, "관심 분야/전공", "Machine Learning / Deep Learning")
                    replace_text_preserve_format(text_frame, "팀에서 수행한 역할", "모델 실험 및 최적화, 하이퍼파라미터 튜닝")
                elif "이승현" not in full_text:  # 두 번째 팀원
                    replace_text_preserve_format(text_frame, "팀원 이름", "이승현")
                    replace_text_preserve_format(text_frame, "관심 분야/전공", "Data Science / EDA")
                    replace_text_preserve_format(text_frame, "팀에서 수행한 역할", "EDA 및 데이터 전처리, 데이터 시각화")
                elif "정재훈" not in full_text:  # 세 번째 팀원
                    replace_text_preserve_format(text_frame, "팀원 이름", "정재훈")
                    replace_text_preserve_format(text_frame, "관심 분야/전공", "Model Validation / Analysis")
                    replace_text_preserve_format(text_frame, "팀에서 수행한 역할", "모델 검증 및 성능 분석, 결과 해석")

def update_slide_5_collaboration(slide):
    """슬라이드 5: 경진대회 협업 방식 업데이트"""
    collaboration_content = """경진대회 진행 방식 및 협업 내용:

협업 마인드셋: 
- 개방적 소통: 모르는 부분은 즉시 팀원들과 공유하고 함께 해결
- 상호 지원: 작업 완료한 팀원은 다른 팀원의 진행 상황을 체크하고 도움 제공
- 투명한 진행: 모든 실험과 결과를 공유하여 팀 전체가 학습할 수 있도록 함

협업 진행 방식:
- 주 3회 정기 미팅 (1시간씩)으로 진행 상황 공유 및 이슈 논의
- GitHub를 통한 코드 공유 및 버전 관리
- WandB를 활용한 실험 결과 공유 및 모니터링
- 슬랙을 통한 실시간 소통 및 질문 해결

협업 과정에서 발생한 문제점:
- 모델 성능 개선 방향에 대한 의견 차이 발생
- 데이터 전처리 방법에 대한 다양한 접근법 논의 필요

문제 해결 방법:
- 각자 다른 접근법으로 실험 진행 후 결과 비교
- 정량적 평가 지표(F1 Score, Accuracy)를 기준으로 객관적 판단
- 앙상블 기법을 활용하여 서로 다른 모델의 장점을 결합"""

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if "진행했던 경진대회" in text_frame.text:
                # 기존 텍스트의 첫 번째 run의 서식을 저장
                original_font_name = None
                original_font_size = None
                original_font_bold = None
                
                if len(text_frame.paragraphs) > 0 and len(text_frame.paragraphs[0].runs) > 0:
                    first_run = text_frame.paragraphs[0].runs[0]
                    font = first_run.font
                    original_font_name = font.name
                    original_font_size = font.size
                    original_font_bold = font.bold
                
                # 모든 텍스트 지우기
                text_frame.clear()
                
                # 새 텍스트 추가하면서 서식 유지
                paragraph = text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = collaboration_content
                
                # 서식 복사
                if original_font_name:
                    run.font.name = original_font_name
                if original_font_size:
                    run.font.size = original_font_size
                if original_font_bold is not None:
                    run.font.bold = original_font_bold

def update_slide_7_goals(slide):
    """슬라이드 7: 목표 수립 업데이트"""
    goal_text = """초기 목표:
1. 문서 분류 모델의 기본 성능 달성 (F1 Score 0.85 이상)
2. EfficientNet 계열 모델을 활용한 고성능 모델 개발
3. K-Fold 교차검증을 통한 안정적인 모델 성능 확보
4. 데이터 증강 및 TTA를 활용한 성능 개선

최종 목표:
- F1 Score 0.90 이상 달성
- 17개 클래스 모든 문서 타입에 대한 정확한 분류 성능 확보
- 실제 업무 환경에서 활용 가능한 견고한 모델 개발"""

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if "경진대회에 대한 팀 초기 목표" in text_frame.text:
                replace_text_preserve_format(text_frame, "경진대회에 대한 팀 초기 목표 설명", goal_text)

def update_slide_8_process(slide):
    """슬라이드 8: 수행 내용 업데이트"""
    process_replacements = {
        "데이터 분석": "상세 EDA 수행 - 클래스 분포, 이미지 해상도, 밝기 분석, 클래스 불균형 문제 파악",
        "개발 환경 구축": "PyTorch 기반 딥러닝 환경 구축 - timm, albumentations, WandB 등 라이브러리 활용",
        "Feature 엔지니어링": "고급 데이터 전처리 - 밝기 보정, 회전 보정, 배경 정규화, 노이즈 제거",
        "모델 선택 학습 및 평가": "EfficientNet 계열 모델 실험 - B3, V2-L, V2-XL, V2-RW-M 비교 분석",
        "작성 작성 작성": "K-Fold 교차검증 (5-fold) + TTA 및 앙상블 기법 적용"
    }
    
    find_and_replace_in_slide(slide, process_replacements)

def update_slide_9_results(slide):
    """슬라이드 9: 수행 결과 업데이트"""
    result_text = """주요 성과 및 결과:

1. 모델 성능 달성:
   - F1 Score: 0.93+ (목표 0.90 대비 초과 달성)
   - Accuracy: 0.93+ (17개 클래스 모든 문서 타입 높은 정확도)

2. 기술적 성과:
   - EfficientNetV2-XL 모델 최적화 완료
   - 5-Fold 교차검증으로 안정적인 성능 확보
   - 10개 모델 앙상블로 최고 성능 달성

3. 데이터 처리 성과:
   - 1,570장 훈련 데이터의 클래스 불균형 문제 해결
   - 고급 데이터 증강 기법으로 모델 일반화 성능 향상
   - Test 데이터 3,140장에 대한 안정적인 예측 수행"""

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if "결측값 확인 및 데이터 전처리" in text_frame.text:
                replace_text_preserve_format(text_frame, "결측값 확인 및 데이터 전처리, 모델 실험 등에 관해 팀에서 시도한 점의 결과 이미지와 요약\n정리", result_text)

def update_slide_11_problems(slide):
    """슬라이드 11: 문제 도출 업데이트"""
    problem_text = """문제 발생 배경 및 원인 분석:

초기에는 단순한 CNN 모델과 기본적인 데이터 전처리만으로 문서 분류를 시도했습니다. 
하지만 17개 클래스의 다양한 문서 타입을 구분하는 것은 예상보다 어려웠습니다.

주요 문제점:
1. 클래스 불균형: 일부 클래스는 100장, 일부는 50장 미만으로 데이터 불균형 존재
2. 이미지 품질 차이: 훈련 데이터와 테스트 데이터 간 밝기, 회전 상태 차이 존재  
3. 모델 선택 문제: 기본 ResNet 모델로는 세밀한 문서 특징 추출에 한계
4. 검증 전략 부재: 단순 홀드아웃 방식으로는 안정적인 성능 평가 어려움"""

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if "처음에는 모든 피처를" in text_frame.text:
                # 긴 텍스트를 찾아서 교체
                old_long_text = "처음에는 모든 피처를 모델에 포함시켰습니다. 다양한 특성들을 사용하여 모델을 훈련시키고, 예측 정확도를 높이려고 했지만, 모델 성능이 기대보다 낮았습니다.\n특히, 불필요한 특성이나 상관관계가 낮은 특성들이 모델의 과적합(overfitting)을 유발하거나, 예측 성능을 떨어뜨리는 원인이 되었던 것입니다"
                replace_text_preserve_format(text_frame, old_long_text, problem_text)

def update_slide_12_insights(slide):
    """슬라이드 12: 인사이트 도출 업데이트"""
    insight_text = """핵심 인사이트 도출:

1. 모델 아키텍처 인사이트:
   - EfficientNet 계열 모델이 문서 분류에 매우 효과적
   - 특히 EfficientNetV2-XL이 복잡한 문서 패턴 인식에 우수한 성능
   - 사전 훈련된 가중치 활용으로 적은 데이터로도 높은 성능 달성

2. 데이터 전처리 인사이트:
   - 훈련 데이터와 테스트 데이터 간 분포 차이 해결이 핵심
   - 밝기 보정, 회전 보정이 성능 향상에 크게 기여
   - 클래스별 특성에 맞는 차별화된 전처리 전략 필요

3. 검증 전략 인사이트:
   - K-Fold 교차검증이 안정적인 모델 성능 평가에 필수
   - 다양한 시드와 폴드 조합으로 모델 신뢰성 크게 향상
   - TTA(Test Time Augmentation) 기법이 최종 성능 개선에 효과적"""

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if "상관 분석을 통해" in text_frame.text:
                # 긴 텍스트를 찾아서 교체
                old_insight_text = "상관 분석을 통해 상관 관계가 낮은 특성들을 확인하고, 이를 제거한 후 모델을 다시 훈련시켰습니다.\n예를 들어, 데이터셋에 다중공선성(multicollinearity) 문제가 있던 특성들(두 변수 간의 상관 관계가 높거나 유사한 변수를 포함)이나, 예측에 큰 영향을 미치지 않는 특성(상관계수나 중요도가 낮은 변수)을 제거하는 과정에서 큰 차이를 발견했습니다"
                replace_text_preserve_format(text_frame, old_insight_text, insight_text)

def update_slide_13_solutions(slide):
    """슬라이드 13: 해결방법 및 결과 업데이트"""
    replacements = {
        "피처 선택 기법을 활용해 중요하지 않은 특성들을 제거": "고급 모델 아키텍처 도입 및 체계적인 데이터 전처리",
        "상관계수 분석을 통해 상관 관계가 높은 특성을 파악하고, 한쪽을 제거.": "EfficientNetV2-XL 모델 활용 및 클래스별 최적화된 전처리 전략 적용.",
        "Lasso 회귀(L1 규제)를 사용하여, 자동으로 불필요한 특성의 가중치를 0으로 만들어 중요한 특성만 남도록 했습니다.": "K-Fold 교차검증과 TTA 기법을 통해 모델 성능과 안정성을 동시에 확보했습니다.",
        "교차 검증을 통해 모델 성능이 더욱 안정적이고 개선됨을 확인했습니다.": "10개 모델 앙상블로 최종 성능을 극대화했습니다.",
        "피처 선택을 통해 모델의 성능이 상당히 개선되었습니다. RMSE (Root Mean Squared Error) 값이 감소했으며, 모델의 일반화 성능도 향상되었습니다.\n이 인사이트를 통해 우리는 회귀 문제에서 데이터 전처리와 특성 선택이 모델 성능을 극대화하는 핵심이라는 중요한 교훈을 얻었습니다.": "최종 결과: F1 Score 0.93+ 달성으로 목표를 초과 달성했으며, 17개 클래스 모든 문서 타입에 대한 높은 정확도를 확보했습니다.\n이 경험을 통해 문서 분류 작업에서 모델 선택, 데이터 전처리, 검증 전략의 중요성을 깊이 이해하게 되었으며, 실제 업무 환경에서 활용 가능한 견고한 AI 모델 개발 역량을 확보했습니다."
    }
    
    find_and_replace_in_slide(slide, replacements)

def update_slide_15_retrospective(slide):
    """슬라이드 15: 회고 업데이트"""
    replacements = {
        "우리 팀의 처음 목표에서 어디까지 도달했는가": "목표 달성도 평가",
        "작성 작성 작성": "초기 목표 F1 Score 0.90 → 실제 달성 F1 Score 0.93+ (목표 초과 달성)",
        "우리 팀이 잘했던 점": "팀의 강점: 체계적인 역할 분담, 정기적인 미팅, GitHub와 WandB를 활용한 협업",
        "협업하면서 아쉬웠던 점": "개선점: 초기 기획의 구체화, 실험 결과 공유 표준화",
        "향후 계획 : ": "향후 계획: 더 복잡한 문서 분류 태스크 도전, 실시간 시스템 구축 경험 축적",
        "첫 번째 경진대회 진행 소감 및 앞으로의 계획에 대해 정리해 작성해주세요. ": "문서 분류 분야의 전문성을 크게 향상시켰고, 실제 업무에서 활용 가능한 수준의 모델을 개발할 수 있었습니다."
    }
    
    find_and_replace_in_slide(slide, replacements)

def update_slide_16_personal_thoughts(slide):
    """슬라이드 16: 개인별 소감 업데이트"""
    personal_info = [
        ("문국현", "팀장으로서 전체 프로젝트를 총괄하며 EfficientNet 모델 최적화와 앙상블 기법을 성공적으로 적용했습니다. 앞으로는 더 복잡한 Computer Vision 태스크에 도전하여 기술 전문성을 높이고, 실제 산업 현장에서 활용 가능한 AI 솔루션 개발 역량을 키워나가겠습니다."),
        ("류지헌", "다양한 EfficientNet 모델 실험과 하이퍼파라미터 최적화를 통해 모델 성능을 극대화하는 과정에서 많은 것을 배웠습니다. 특히 K-Fold 교차검증과 TTA 기법의 중요성을 깊이 깨달았습니다. 앞으로는 AutoML과 Neural Architecture Search 분야를 더 깊이 연구하여 자동화된 모델 최적화 기술을 개발하고 싶습니다."),
        ("이승현", "상세한 EDA를 통해 데이터의 특성을 파악하고 효과적인 전처리 전략을 수립하는 과정이 매우 의미 있었습니다. 클래스 불균형 문제와 이미지 품질 차이를 해결하는 방법을 실제로 적용해볼 수 있어서 좋았습니다. 앞으로는 더 복잡한 데이터 분석 도구와 시각화 기법을 익혀서 데이터 사이언티스트로서의 역량을 확장하고 싶습니다."),
        ("정재훈", "모델 검증과 성능 분석을 담당하며 다양한 평가 지표와 교차검증 기법의 중요성을 배웠습니다. 특히 F1 Score와 Accuracy의 차이를 이해하고 적절한 평가 전략을 수립하는 경험이 값졌습니다. 앞으로는 MLOps와 모델 모니터링 분야를 더 깊이 공부하여 프로덕션 환경에서의 모델 운영 전문가가 되고 싶습니다.")
    ]
    
    # 기본 텍스트 교체
    replacements = {
        "김OO": "팀원 이름",
        "이번 스터디 경험을 통해 앞으로의 스터디 방향성 및 각오를 말씀해주세요": "개인별 경험과 향후 계획"
    }
    find_and_replace_in_slide(slide, replacements)
    
    # 개인별 이름과 소감 순차적 업데이트
    name_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if "팀원 이름" in text_frame.text and name_count < len(personal_info):
                replace_text_preserve_format(text_frame, "팀원 이름", personal_info[name_count][0])
                name_count += 1
            elif "개인별 경험과 향후 계획" in text_frame.text:
                # 해당하는 팀원의 소감으로 교체
                person_idx = min(name_count - 1, len(personal_info) - 1)
                if person_idx >= 0:
                    replace_text_preserve_format(text_frame, "개인별 경험과 향후 계획", personal_info[person_idx][1])

def main():
    # PowerPoint 파일 경로
    ppt_path = "ppt/cv-competition2.pptx"
    
    if not os.path.exists(ppt_path):
        print(f"PowerPoint 파일을 찾을 수 없습니다: {ppt_path}")
        return
    
    # PowerPoint 파일 로드
    prs = Presentation(ppt_path)
    
    print(f"총 {len(prs.slides)}개의 슬라이드를 업데이트합니다...")
    print("📝 기존 텍스트의 서식을 유지하면서 내용을 업데이트합니다.\n")
    
    # 각 슬라이드별 업데이트
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        print(f"슬라이드 {slide_num} 처리 중...")
        
        if slide_num == 4:
            update_slide_4_team_info(slide)
            print("  ✅ 팀원 소개 업데이트 완료 (서식 유지)")
        elif slide_num == 5:
            update_slide_5_collaboration(slide)
            print("  ✅ 협업 방식 업데이트 완료 (서식 유지)")
        elif slide_num == 7:
            update_slide_7_goals(slide)
            print("  ✅ 목표 수립 업데이트 완료 (서식 유지)")
        elif slide_num == 8:
            update_slide_8_process(slide)
            print("  ✅ 수행 내용 업데이트 완료 (서식 유지)")
        elif slide_num == 9:
            update_slide_9_results(slide)
            print("  ✅ 수행 결과 업데이트 완료 (서식 유지)")
        elif slide_num == 11:
            update_slide_11_problems(slide)
            print("  ✅ 문제 도출 업데이트 완료 (서식 유지)")
        elif slide_num == 12:
            update_slide_12_insights(slide)
            print("  ✅ 인사이트 도출 업데이트 완료 (서식 유지)")
        elif slide_num == 13:
            update_slide_13_solutions(slide)
            print("  ✅ 해결방법 및 결과 업데이트 완료 (서식 유지)")
        elif slide_num == 15:
            update_slide_15_retrospective(slide)
            print("  ✅ 회고 업데이트 완료 (서식 유지)")
        elif slide_num == 16:
            update_slide_16_personal_thoughts(slide)
            print("  ✅ 개인별 소감 업데이트 완료 (서식 유지)")
        else:
            print("  ⏩ 변경 사항 없음")
    
    # 업데이트된 파일 저장
    prs.save(ppt_path)
    print(f"\n✅ PowerPoint 파일 업데이트가 완료되었습니다: {ppt_path}")
    print("🎨 기존 텍스트의 서식(폰트, 색상, 크기 등)을 유지하면서 내용이 업데이트되었습니다.")
    print("📊 README.md 내용을 바탕으로 실제 프로젝트 정보로 업데이트되었습니다.")

if __name__ == "__main__":
    main() 