
from pptx import Presentation

def inspect_presentation(pptx_path):
    """PPTX 파일의 각 슬라이드에 있는 도형의 텍스트를 출력합니다."""
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(f"  Shape Text: {shape.text_frame.text}")
            else:
                print("  (No text in this shape)")

if __name__ == '__main__':
    inspect_presentation('cv-competition_backup.pptx')
